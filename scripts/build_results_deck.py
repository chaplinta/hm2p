#!/usr/bin/env python3
"""Build a slide deck summarising the RSP light/dark neural + behaviour results.

Produces docs/manuscripts/results-summary-deck.pptx. Numbers are taken from the
committed analyses (results/dark_hypotheses, results/map_engagement, and the
behaviour summary). 16:9, neutral scientific styling.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "docs" / "manuscripts" / "results-summary-deck.pptx"

DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
MID_BLUE = RGBColor(0x2E, 0x75, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MID_GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x2D, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _deck() -> Presentation:
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title_bar(slide, text):
    box = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BLUE
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.5)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = WHITE


def _bullets(slide, items, top=1.3, left=0.6, width=12.1, height=5.8, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level, colour, bold = item
        else:
            text, level, colour, bold = item, 0, DARK_GREY, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = ("• " if level == 0 else "– ") + text
        r.font.size = Pt(size - 2 * level)
        r.font.color.rgb = colour
        r.font.bold = bold


def _table(slide, headers, rows, top=1.4, left=0.6, width=12.1, col_w=None, size=14):
    nr, nc = len(rows) + 1, len(headers)
    gfx = slide.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(width), Inches(0.4 * nr))
    tbl = gfx.table
    if col_w:
        for j, w in enumerate(col_w):
            tbl.columns[j].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = MID_BLUE
        c.text = h
        for para in c.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(size)
                run.font.bold = True
                run.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            text, colour = val if isinstance(val, tuple) else (val, DARK_GREY)
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT_GREY
            c.text = text
            for para in c.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(size)
                    run.font.color.rgb = colour


def _footnote(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.italic = True
    r.font.color.rgb = MID_GREY


def build():
    prs = _deck()

    # 1. Title
    s = _blank(prs)
    bar = s.shapes.add_shape(1, 0, Inches(2.4), SLIDE_W, Inches(2.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Head-direction coding and exploration in retrosplenial cortex under light and dark"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Freely-moving 2P imaging, q-rose maze — results summary"
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)

    # 2. Setup
    s = _blank(prs)
    _title_bar(s, "Setup")
    _bullets(s, [
        "Two-photon GCaMP imaging (~9.6 Hz, single plane) of retrosplenial cortex in freely-moving mice",
        "Two non-overlapping RSP populations: Penk+ and Penk-CamKII+ (12 vs 4 animals)",
        "q-rose maze; overhead lights cycle 1 min on / 1 min off",
        ("Light off = total darkness = visual cues removed. Camera is infrared, so tracking is unaffected — only the mouse's vision changes", 1, MID_GREY, False),
        "Within-session light vs dark, paired. 23 sessions, 15 animals. Non-parametric throughout",
        ("Question: how does removing vision change head-direction coding and spatial behaviour?", 0, DARK_BLUE, True),
    ])

    # 3. Starting observation
    s = _blank(prs)
    _title_bar(s, "Starting observation: a counterintuitive effect")
    _bullets(s, [
        ("Raw within-cell result: HD tuning (mean vector length) and spatial information are HIGHER in dark than light", 0, GREEN, True),
        ("dF/F MVL p = 0.0010; spatial info p < 0.0001; preferred directions show no systematic drift", 1, DARK_GREY, False),
        "This runs opposite to the classical expectation that darkness degrades head-direction tuning",
        "~9.8% of soma ROIs are HD-tuned (dF/F)",
        ("Too good to take at face value: MVL is biased by how the animal samples head direction", 0, RED, True),
    ])

    # 4. The confound
    s = _blank(prs)
    _title_bar(s, "Why the raw effect is suspect")
    _bullets(s, [
        "MVL of a tuning curve rises when head-direction sampling is narrower or more uneven, with no change in the underlying code",
        "In the dark the mouse samples head direction differently (it explores a smaller, more repeated set of routes)",
        "So a higher dark MVL could be a sampling artefact, not a coding gain",
        ("Test: equalise the behavioural sampling between light and dark, then recompute", 0, DARK_BLUE, True),
        ("A1 — match the head-direction occupancy distribution", 1, DARK_GREY, False),
        ("A2 — match the joint speed x |angular head velocity| distribution", 1, DARK_GREY, False),
        ("Plus circular-shuffle debiasing of MVL; within-session paired", 1, DARK_GREY, False),
    ])

    # 5. Gauntlet result
    s = _blank(prs)
    _title_bar(s, "The effect does not survive matching")
    _table(s,
           ["Control", "Light vs dark (MVL)", "Verdict"],
           [
               ["Raw (no matching)", "p = 0.0010", ("dark > light", GREEN)],
               ["A1 occupancy-matched", "p = 0.16", ("not significant", RED)],
               ["A2 speed+|AHV|-matched", "p = 0.25", ("not significant", RED)],
               ["A3 epoch-order / bleaching", "p = 0.10", ("not significant", RED)],
           ],
           col_w=[5.0, 4.0, 3.1])
    _bullets(s, [
        "Equalising how the mouse samples head direction removes the effect",
        ("The raw dark>light MVL is largely differential sampling, not a coding gain", 0, DARK_BLUE, True),
    ], top=4.6)
    _footnote(s, "Sanity checks pass: matching disabled reproduces raw MVL exactly; occupancy matching equalises the HD histograms.")

    # 6. MI cross-check
    s = _blank(prs)
    _title_bar(s, "A second metric agrees: mutual information")
    _bullets(s, [
        "Skaggs head-direction information (bits/event; Voigts & Harnett 2020, Zong et al. 2022) run through the same matched gauntlet",
        ("Matched HD-MI is null, more emphatically than MVL: A1 p = 0.64, A2 p = 0.89", 0, DARK_GREY, False),
        "Place (spatial) information also dies under position-occupancy matching (p = 0.15)",
        "Secondary signals collapse too: the apparent 'gain' does not survive (matched p = 0.21); the apparent dark cell-recruitment reverses to favour light",
        ("Two independent metrics give the same answer: no dark enhancement of HD or place coding survives matching", 0, DARK_BLUE, True),
    ])

    # 7. Neural conclusion
    s = _blank(prs)
    _title_bar(s, "Neural conclusion")
    _bullets(s, [
        ("The RSP spatial / head-direction representation is preserved in the dark — neither enhanced nor degraded", 0, GREEN, True),
        "Every matched measure is null: HD MVL, HD mutual information, place information, and population-vector map consistency",
        "The naive 'darkness enhances HD coding' was a sampling artefact and is retired",
        "Cell-type contrast (Penk+ vs Penk-CamKII+) is null and underpowered: 12 vs 4 animals can only detect very large effects (d > 1.5)",
        ("Controls held: tracking confidence is identical light vs dark (infrared camera assumption upheld)", 1, MID_GREY, False),
    ])

    # 8. Behaviour: same locomotion
    s = _blank(prs)
    _title_bar(s, "Behaviour: locomotion is unchanged")
    _table(s,
           ["Measure", "Light", "Dark", "Test"],
           [
               ["Speed (cm/s)", "2.41", "2.12", ("p = 0.07 (ns)", MID_GREY)],
               ["Distance/epoch (m)", "2.26", "2.03", ("p = 0.03 raw, ns adj", MID_GREY)],
               ["Left-turn fraction", "0.50", "0.50", ("p = 0.99", MID_GREY)],
               ["Turn alternation", "-0.16", "-0.18", ("p = 0.54", MID_GREY)],
               ["Head angular velocity", "same", "same", ("p = 0.18", MID_GREY)],
           ],
           col_w=[4.6, 2.3, 2.3, 3.0])
    _bullets(s, [
        "Same speed, same distance, same turn choices, same head movement — the motor machinery is intact in the dark",
    ], top=4.9)

    # 9. Behaviour: exploration changes
    s = _blank(prs)
    _title_bar(s, "Behaviour: exploration changes")
    _table(s,
           ["Measure", "Light", "Dark", "Corrected p"],
           [
               ["Coverage (unique cells/min)", "0.41", "0.35", ("0.001", GREEN)],
               ["Occupancy entropy (bits)", "2.24", "1.96", ("0.009", GREEN)],
               ["Coverage vs random-walk null (z)", "0.69", "0.23", ("0.018", GREEN)],
               ["Route compressibility (LZ)", "0.73", "0.72", ("0.85 (ns)", MID_GREY)],
           ],
           col_w=[5.4, 2.0, 2.0, 2.7])
    _bullets(s, [
        ("In light, mice cover more of the maze than a random walk of the same length — directed search", 0, DARK_BLUE, True),
        ("In dark they fall to near random-walk level, controlling for how much they moved and the maze shape", 0, DARK_BLUE, True),
        ("Local route patterns are unchanged (LZ) — they don't run tighter loops; they stop steering toward new ground", 1, MID_GREY, False),
    ], top=4.0)

    # 10. Map engagement
    s = _blank(prs)
    _title_bar(s, "Is the map still engaged in the dark? Yes")
    _bullets(s, [
        "Test: does the same maze cell re-instantiate the same population state on each visit? (within- minus across-cell population-vector consistency, sampling matched, no decoding)",
        ("Consistency is positive in nearly every session — a real spatial map exists", 0, GREEN, True),
        ("It does NOT differ light vs dark: 0.065 vs 0.062, p = 0.49 (12/21 sessions light>dark)", 0, DARK_BLUE, True),
        "The spatial representation is equally engaged with or without vision",
        ("Caveats: small/noisy measure (~20 ROIs, underpowered to prove equivalence); HD cells not excluded", 1, MID_GREY, False),
    ])

    # 11. Synthesis
    s = _blank(prs)
    _title_bar(s, "Synthesis")
    _bullets(s, [
        ("The map is intact in the dark; the animal stops using it to direct exploration", 0, DARK_BLUE, True),
        "Neural side: HD, place, and map consistency are all preserved without vision",
        "Behaviour side: same locomotion, but exploration goes from directed to random-walk-like",
        "The dissociation is between representation (intact) and its behavioural readout (changes) — not a change in the code itself",
        ("This is a cleaner story than either 'darkness degrades' or 'darkness enhances' coding, both of which were ruled out", 0, GREEN, True),
    ])

    # 12. Caveats & next
    s = _blank(prs)
    _title_bar(s, "Caveats and next steps")
    _bullets(s, [
        "HD yield is ~9.8% (2-4 HD cells per session) — population decoding may not be viable; this bounds the neural claims",
        "Matched nulls are 'no detectable difference', not proven equivalence — would firm up with more shuffles / Bayesian framing",
        "Cell-type comparisons need more Penk-CamKII+ animals to be informative",
        ("Next, if pursued:", 0, DARK_BLUE, True),
        ("AHV / egocentric coding light vs dark (predicted preserved; currently unrun)", 1, DARK_GREY, False),
        ("Junction-choice predictability and per-epoch neural-behaviour coupling", 1, DARK_GREY, False),
        ("HD-excluded map-engagement and a symmetric-corridor version of the junction analysis", 1, DARK_GREY, False),
    ])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()

#!/usr/bin/env python3
"""Generate a PowerPoint presentation from the hm2p architecture slide deck.

Reads ``docs/architecture-slides.md`` and produces
``docs/architecture-presentation.pptx`` with shape-based diagrams, tables,
and formatted bullet-point slides.

Usage::

    python3 scripts/build_architecture_pptx.py
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
MID_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MID_GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT_GREEN = RGBColor(0x2D, 0x8B, 0x57)
ACCENT_ORANGE = RGBColor(0xE0, 0x7C, 0x24)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_TEAL = RGBColor(0x17, 0x85, 0x85)

# ---------------------------------------------------------------------------
# Slide dimensions (widescreen 16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Font helpers
# ---------------------------------------------------------------------------

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
FONT_MONO = "Consolas"


def _set_font(
    run,
    *,
    size: int = 18,
    bold: bool = False,
    italic: bool = False,
    colour: RGBColor = DARK_GREY,
    name: str = FONT_BODY,
):
    """Apply font formatting to a *run*."""
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour


def _add_run(
    paragraph,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    italic: bool = False,
    colour: RGBColor = DARK_GREY,
    name: str = FONT_BODY,
):
    """Append a formatted run to *paragraph* and return it."""
    run = paragraph.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, italic=italic, colour=colour, name=name)
    return run


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------


def _add_title_bar(slide, title_text: str):
    """Add a dark-blue title bar across the top of the slide."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(0),
        Emu(0),
        SLIDE_W,
        Inches(1.0),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    _set_font(run, size=28, bold=True, colour=WHITE, name=FONT_TITLE)
    # Vertical centering
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)


def _add_subtitle(slide, text: str, top: float = 1.15):
    """Add a subtitle line below the title bar."""
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.4))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    _add_run(p, text, size=14, italic=True, colour=MID_GREY)


def _add_notes(slide, text: str):
    """Set speaker notes on *slide*."""
    slide.notes_slide.notes_text_frame.text = text


def _add_bullet_box(
    slide,
    bullets: list[str],
    left: float = 0.5,
    top: float = 1.5,
    width: float = 12.3,
    height: float = 5.5,
    font_size: int = 16,
    bold_before_colon: bool = True,
):
    """Add a text box with bullet points. Bold text before the first colon."""
    txbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txbox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        p.level = 0

        # Determine indent level from leading whitespace / dash
        stripped = bullet.lstrip()
        if bullet.startswith("  ") or bullet.startswith("    "):
            p.level = 1

        # Clean leading bullets/dashes
        text = re.sub(r"^[-*]\s*", "", stripped)

        # Bold before first colon if present
        if bold_before_colon and ":" in text and not text.startswith("("):
            colon_pos = text.index(":")
            bold_part = text[: colon_pos + 1]
            rest = text[colon_pos + 1 :]
            _add_run(p, bold_part, size=font_size, bold=True, colour=DARK_BLUE)
            if rest:
                _add_run(p, rest, size=font_size, colour=DARK_GREY)
        else:
            _add_run(p, text, size=font_size, colour=DARK_GREY)


def _rect(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    fill: RGBColor = MID_BLUE,
    font_colour: RGBColor = WHITE,
    font_size: int = 12,
    bold: bool = True,
):
    """Draw a rounded-rectangle box with centred text."""
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        Emu(int(left)),
        Emu(int(top)),
        Emu(int(width)),
        Emu(int(height)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = DARK_BLUE
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    # Split text on newlines for multi-line labels
    lines = text.split("\n")
    for j, line in enumerate(lines):
        if j > 0:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        _set_font(run, size=font_size, bold=bold, colour=font_colour)

    return shape


def _arrow_right(slide, x1, y, x2, colour: RGBColor = DARK_BLUE):
    """Draw a right-pointing arrow connector from (x1, y) to (x2, y)."""
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Emu(int(x1)),
        Emu(int(y)),
        Emu(int(x2)),
        Emu(int(y)),
    )
    connector.line.color.rgb = colour
    connector.line.width = Pt(2)
    # Arrowhead
    connector.line.end_marker_style = 1  # arrow
    return connector


def _arrow_down(slide, x, y1, y2, colour: RGBColor = DARK_BLUE):
    """Draw a downward arrow from (x, y1) to (x, y2)."""
    connector = slide.shapes.add_connector(
        1,
        Emu(int(x)),
        Emu(int(y1)),
        Emu(int(x)),
        Emu(int(y2)),
    )
    connector.line.color.rgb = colour
    connector.line.width = Pt(2)
    connector.line.end_marker_style = 1
    return connector


# ---------------------------------------------------------------------------
# Table helper
# ---------------------------------------------------------------------------


def _add_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    left: float = 0.5,
    top: float = 1.5,
    width: float = 12.3,
    row_height: float = 0.35,
    header_font_size: int = 13,
    body_font_size: int = 12,
    col_widths: list[float] | None = None,
):
    """Add a formatted table to the slide."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_height = row_height * n_rows

    table_shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(tbl_height),
    )
    table = table_shape.table

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        _add_run(p, h, size=header_font_size, bold=True, colour=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            _add_run(p, val, size=body_font_size, colour=DARK_GREY)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 0 else LIGHT_GREY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


# ---------------------------------------------------------------------------
# Parse speaker notes from markdown
# ---------------------------------------------------------------------------


def _parse_slides(md_path: Path) -> list[dict]:
    """Parse the markdown into a list of slide dicts with title, body, notes."""
    text = md_path.read_text()
    # Split on slide headers
    slide_blocks = re.split(r"^## Slide \d+:\s*", text, flags=re.MULTILINE)
    # First block is the preamble (before Slide 1)
    slides = []
    for block in slide_blocks[1:]:
        lines = block.strip().split("\n")
        title = lines[0].strip()
        # Find speaker notes
        body_lines = []
        notes_lines = []
        in_notes = False
        for line in lines[1:]:
            if line.strip().startswith("**Speaker notes:**"):
                in_notes = True
                # The rest of this line might have content
                after = line.replace("**Speaker notes:**", "").strip()
                if after:
                    notes_lines.append(after)
                continue
            if in_notes:
                notes_lines.append(line)
            else:
                body_lines.append(line)
        slides.append(
            {
                "title": title,
                "body": "\n".join(body_lines).strip(),
                "notes": "\n".join(notes_lines).strip(),
            }
        )
    return slides


# ---------------------------------------------------------------------------
# Individual slide builders
# ---------------------------------------------------------------------------


def _build_title_slide(prs: Presentation, slide_data: dict):
    """Slide 1: The Experiment (title-style with bullets)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_title_bar(slide, "The Experiment")
    _add_subtitle(slide, "hm2p-v2 Architecture Overview")

    bullets = [
        "Animal: Freely-moving mouse in Rosenberg q-rose maze (7x5 grid, 23 cells, 6 dead-ends)",
        "Brain region: Retrosplenial cortex (RSP) -- head-direction (HD) cells",
        "Imaging: Single-plane 2P GCaMP7f at ~9.6 Hz (SciScan, 512x512 px)",
        "Behaviour: Overhead video at ~100 fps (Basler acA1300-200um), DAQ-synchronised",
        "Light manipulation: 1 min on / 1 min off -- total darkness = visual cue removal",
        "Two non-overlapping RSP populations:",
        "  Penk+: Penk-Cre + AAV-ADD3 (Cre-ON)",
        "  Penk-minus CamKII+: Penk-Cre + virus 344 (Cre-OFF intersectional)",
        "Dataset: 26 sessions across 9 animals; ~113 GB raw data",
    ]
    _add_bullet_box(slide, bullets, top=1.55, font_size=15)
    _add_notes(slide, slide_data["notes"])


def _build_raw_data_slide(prs: Presentation, slide_data: dict):
    """Slide 2: Raw Data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Raw Data -- What Comes Off the Rig")

    bullets = [
        "Two-photon TIFFs: .raw (SciScan) converted to .tif stacks; green + red channels",
        "DAQ timing file: .tdms with trigger pulses (camera, SciScan, light on/off)",
        "Overhead video: .mp4 (H.264), pre-processed (undistorted + cropped)",
        "Video metadata: meta/meta.txt -- crop region, scale (mm/px), maze ROI corners",
        "SciScan metadata: .meta.txt -- frame rate, DAQ channel map",
        "Camera calibration: lens-specific .npz files (4 mm and 6 mm)",
        "Z-stacks: serial2p multi-page TIFFs (16/26 sessions, 13 z-stacks)",
        "Whole-brain volumes: post-mortem serial2p (25 um), registered to Allen CCFv3",
    ]
    _add_bullet_box(slide, bullets, top=1.3, font_size=15)
    _add_notes(slide, slide_data["notes"])


def _build_data_flow_slide(prs: Presentation, slide_data: dict):
    """Slide 3: Data Flow -- pipeline diagram with shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Data Flow -- From Raw to Analysis-Ready")

    # Dimensions in EMU (English Metric Units)
    inch = 914400  # 1 inch in EMU

    # Stage boxes -- arranged as a flow diagram
    box_w = int(2.0 * inch)
    box_h = int(0.7 * inch)
    small_box_h = int(0.55 * inch)

    # Row 1: Raw Data
    raw_x = int(5.5 * inch)
    raw_y = int(1.3 * inch)
    _rect(slide, raw_x, raw_y, box_w, small_box_h, "Raw Data\n(Dropbox/S3)", fill=DARK_GREY)

    # Row 2: Stage 0
    s0_x = raw_x
    s0_y = int(2.2 * inch)
    _rect(slide, s0_x, s0_y, box_w, small_box_h, "Stage 0\nIngest & DAQ Parse", fill=DARK_BLUE)
    _arrow_down(slide, raw_x + box_w // 2, raw_y + small_box_h, s0_y)

    # Row 3: Branch -- Stage 1 (left) and Stage 2a (right)
    s1_x = int(2.5 * inch)
    s2a_x = int(8.5 * inch)
    branch_y = int(3.1 * inch)

    _rect(
        slide,
        s1_x,
        branch_y,
        box_w,
        box_h,
        "Stage 1\n2P Extraction\n(Suite2p + Cellpose 3)",
        fill=MID_BLUE,
    )
    _rect(
        slide,
        s2a_x,
        branch_y,
        box_w,
        box_h,
        "Stage 2a\nDLC Training\n(SuperAnimal fine-tune)",
        fill=ACCENT_TEAL,
    )

    # Arrows from Stage 0 to branches
    _arrow_down(slide, s0_x + box_w // 4, s0_y + small_box_h, branch_y, colour=MID_BLUE)
    _arrow_down(slide, s0_x + 3 * box_w // 4, s0_y + small_box_h, branch_y, colour=ACCENT_TEAL)

    # Horizontal guide lines to branches
    # Left branch connector
    connector_l = slide.shapes.add_connector(
        1,
        Emu(int(s0_x + box_w // 4)),
        Emu(int(s0_y + small_box_h + int(0.25 * inch))),
        Emu(int(s1_x + box_w // 2)),
        Emu(int(s0_y + small_box_h + int(0.25 * inch))),
    )
    connector_l.line.color.rgb = MID_BLUE
    connector_l.line.width = Pt(2)

    # Right branch connector
    connector_r = slide.shapes.add_connector(
        1,
        Emu(int(s0_x + 3 * box_w // 4)),
        Emu(int(s0_y + small_box_h + int(0.25 * inch))),
        Emu(int(s2a_x + box_w // 2)),
        Emu(int(s0_y + small_box_h + int(0.25 * inch))),
    )
    connector_r.line.color.rgb = ACCENT_TEAL
    connector_r.line.width = Pt(2)

    # Vertical drops from horizontal guides
    _arrow_down(
        slide, s1_x + box_w // 2, s0_y + small_box_h + int(0.25 * inch), branch_y, colour=MID_BLUE
    )
    _arrow_down(
        slide,
        s2a_x + box_w // 2,
        s0_y + small_box_h + int(0.25 * inch),
        branch_y,
        colour=ACCENT_TEAL,
    )

    # Row 4: Stage 2b (right side)
    s2b_y = int(4.1 * inch)
    _rect(slide, s2a_x, s2b_y, box_w, small_box_h, "Stage 2b\nDLC Inference", fill=ACCENT_TEAL)
    _arrow_down(slide, s2a_x + box_w // 2, branch_y + box_h, s2b_y, colour=ACCENT_TEAL)

    # Row 4 left: Stage 4
    s4_y = int(4.1 * inch)
    _rect(
        slide,
        s1_x,
        s4_y,
        box_w,
        box_h,
        "Stage 4\nCalcium Processing\n(dF/F, CASCADE, events)",
        fill=MID_BLUE,
    )
    _arrow_down(slide, s1_x + box_w // 2, branch_y + box_h, s4_y, colour=MID_BLUE)

    # Row 5 right: Stage 3
    s3_y = int(4.95 * inch)
    _rect(
        slide,
        s2a_x,
        s3_y,
        box_w,
        box_h,
        "Stage 3\nKinematics\n(HD, position, speed)",
        fill=ACCENT_TEAL,
    )
    _arrow_down(slide, s2a_x + box_w // 2, s2b_y + small_box_h, s3_y, colour=ACCENT_TEAL)

    # Row 6: Convergence -- Stage 5 (center)
    s5_x = int(5.5 * inch)
    s5_y = int(6.0 * inch)
    _rect(
        slide,
        s5_x,
        s5_y,
        box_w,
        small_box_h,
        "Stage 5 -- Sync\n(neural + behavioural merge)",
        fill=DARK_BLUE,
    )

    # Arrows converging to Stage 5
    _arrow_down(slide, s1_x + box_w // 2, s4_y + box_h, s5_y + small_box_h // 2, colour=MID_BLUE)
    # Horizontal from left to center
    slide.shapes.add_connector(
        1,
        Emu(int(s1_x + box_w // 2)),
        Emu(int(s5_y + small_box_h // 2)),
        Emu(int(s5_x)),
        Emu(int(s5_y + small_box_h // 2)),
    ).line.color.rgb = MID_BLUE

    _arrow_down(
        slide, s2a_x + box_w // 2, s3_y + box_h, s5_y + small_box_h // 2, colour=ACCENT_TEAL
    )
    # Horizontal from right to center
    conn = slide.shapes.add_connector(
        1,
        Emu(int(s2a_x + box_w // 2)),
        Emu(int(s5_y + small_box_h // 2)),
        Emu(int(s5_x + box_w)),
        Emu(int(s5_y + small_box_h // 2)),
    )
    conn.line.color.rgb = ACCENT_TEAL

    # Stage 6 arrow (right of Stage 5)
    s6_x = int(8.5 * inch)
    _rect(slide, s6_x, s5_y, box_w, small_box_h, "Stage 6\nAnalysis", fill=ACCENT_ORANGE)
    _arrow_right(slide, s5_x + box_w, s5_y + small_box_h // 2, s6_x, colour=ACCENT_ORANGE)

    # Legend labels
    for lbl, colour, y_off in [
        ("Calcium branch", MID_BLUE, 1.3),
        ("Pose branch", ACCENT_TEAL, 1.6),
        ("Convergence", DARK_BLUE, 1.9),
    ]:
        dot = slide.shapes.add_shape(1, Inches(0.4), Inches(y_off), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = colour
        dot.line.fill.background()
        lbl_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(y_off - 0.02), Inches(1.8), Inches(0.25)
        )
        p = lbl_box.text_frame.paragraphs[0]
        _add_run(p, lbl, size=11, colour=colour, bold=True)

    _add_notes(slide, slide_data["notes"])


def _build_s3_layout_slide(prs: Presentation, slide_data: dict):
    """Slide 4: S3 Storage Layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "S3 Storage Layout")

    # Draw two bucket boxes
    inch = 914400

    # Left bucket: hm2p-rawdata
    _rect(
        slide,
        int(0.5 * inch),
        int(1.4 * inch),
        int(5.8 * inch),
        int(4.3 * inch),
        "",
        fill=LIGHT_BLUE,
        font_colour=DARK_BLUE,
    )
    # Bucket title
    txb = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.4), Inches(0.4))
    p = txb.text_frame.paragraphs[0]
    _add_run(p, "s3://hm2p-rawdata/", size=16, bold=True, colour=DARK_BLUE, name=FONT_MONO)

    raw_items = [
        "rawdata/sub-{id}/ses-{ts}/",
        "  funcimg/  --  2P TIFF stacks + .meta.txt",
        "  behav/  --  .mp4 video + meta/ + .tdms",
        "",
        "sourcedata/",
        "  trackers/dlc/  --  DLC model weights + labeled data",
        "  calibration/  --  camera .npz files",
        "  metadata/  --  animals.csv, experiments.csv",
        "  zstacks/{id}/  --  serial2p z-stacks",
    ]
    txb2 = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(5.4), Inches(3.5))
    tf = txb2.text_frame
    tf.word_wrap = True
    for i, line in enumerate(raw_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _add_run(p, line, size=12, colour=DARK_GREY, name=FONT_MONO)

    # Right bucket: hm2p-derivatives
    _rect(
        slide,
        int(6.7 * inch),
        int(1.4 * inch),
        int(6.1 * inch),
        int(4.3 * inch),
        "",
        fill=LIGHT_BLUE,
        font_colour=DARK_BLUE,
    )
    txb3 = slide.shapes.add_textbox(Inches(6.9), Inches(1.5), Inches(5.7), Inches(0.4))
    p = txb3.text_frame.paragraphs[0]
    _add_run(p, "s3://hm2p-derivatives/", size=16, bold=True, colour=DARK_BLUE, name=FONT_MONO)

    deriv_items = [
        "derivatives/",
        "  ca_extraction/{sub}/{ses}/suite2p/",
        "  pose/{sub}/{ses}/  --  DLC .h5 files",
        "  movement/{sub}/{ses}/kinematics.h5",
        "  calcium/{sub}/{ses}/ca.h5",
        "  sync/{sub}/{ses}/sync.h5",
        "  analysis/{sub}/{ses}/analysis.h5",
        "",
        "dlc-retrain/  --  training data + models",
        "dlc-champion.json  --  current champion model",
        "dlc-champion-history/  --  superseded champions",
    ]
    txb4 = slide.shapes.add_textbox(Inches(6.9), Inches(2.0), Inches(5.7), Inches(3.5))
    tf2 = txb4.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(deriv_items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        _add_run(p, line, size=12, colour=DARK_GREY, name=FONT_MONO)

    # Bottom notes
    bottom_bullets = [
        "NeuroBlueprint folder standard (BIDS-inspired, from SWC/UCL neuroinformatics)",
        "Session names: ses-{YYYYMMDD}T{HHMMSS} (full timestamp; multiple sessions/day)",
        "Two buckets: rawdata (Infrequent Access) vs derivatives (Standard)",
        "Storage cost: ~$10/month (600 GB raw + ~150 GB derivatives)",
    ]
    _add_bullet_box(slide, bottom_bullets, top=5.9, height=1.5, font_size=13)
    _add_notes(slide, slide_data["notes"])


def _build_daq_slide(prs: Presentation, slide_data: dict):
    """Slide 5: Stage 0 - DAQ Parsing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Stage 0 -- DAQ Parsing")

    bullets = [
        "Input: daq.tdms per session (NI DAQ binary format)",
        "Parser: nptdms library -- isolated to src/hm2p/ingest/daq.py",
        "Output: timestamps.h5 with:",
        "  frame_times_camera (N,) float64 -- camera trigger timestamps",
        "  frame_times_imaging (T,) float64 -- 2P frame timestamps",
        "  light_on_times / light_off_times (L,) float64 -- lighting pulse edges",
        "  fps_camera, fps_imaging attributes",
        "Validation: ingest/validate.py checks raw file completeness",
        "Runner: scripts/run_stage0_daq.py",
        "Status: 26/26 sessions complete",
    ]
    _add_bullet_box(slide, bullets, font_size=15)
    _add_notes(slide, slide_data["notes"])


def _build_extraction_slide(prs: Presentation, slide_data: dict):
    """Slide 6: Stage 1 - Two-Photon Extraction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Stage 1 -- Two-Photon Extraction")

    bullets = [
        "Default extractor: Suite2p (Cellpose 3 anatomical prior, anatomical_only=2)",
        "Alternative: CaImAn (CNMF-based) in extraction/caiman.py",
        "Unified API: roiextractors -- seg.get_traces('raw'), get_roi_image_masks()",
        "Cellpose 3: seeds ROI candidates from projections, biases toward compact somata",
        "Soma vs dendrite classification: post-hoc from stat.npy shape statistics",
        "  extraction/soma_features.py -- per-ROI feature extraction",
        "  extraction/soma_classifier.py -- rule-based scorer or logistic regression",
        "  Output: roi_types (0=soma, 1=dendrite, 2=artefact) + calibrated probabilities",
        "Manual ROI curation: extraction/curation.py -- append-only label CSV",
        "Z-drift: extraction/zdrift.py -- register imaging frames against z-stacks",
        "Status: 26/26 sessions processed on EC2 g4dn.xlarge",
    ]
    _add_bullet_box(slide, bullets, font_size=14)
    _add_notes(slide, slide_data["notes"])


def _build_dlc_slide(prs: Presentation, slide_data: dict):
    """Slide 7: Stage 2 - DLC Pose Tracking."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Stage 2 -- DLC Pose Tracking")

    bullets = [
        "Model: DLC 3.0 (PyTorch) + SuperAnimal TopViewMouse + HRNet-W32 + FasterRCNN",
        "8 bodyparts: nose_tip, left_ear, right_ear, head_midpoint, neck, "
        "mid_back, mouse_center, tail_base",
        "  head_midpoint is a custom keypoint (high-contrast 2P headstage)",
        "Stage 2a -- Training (GPU, 24h max):",
        "  Fine-tune SuperAnimal on 183 manually labeled frames",
        "  EC2 launch: scripts/launch_dlc_finetune_ec2.py",
        "  W&B logging for training metrics",
        "Stage 2b -- Inference (GPU):",
        "  deeplabcut.analyze_videos() on all 26 sessions",
        "  Output: .h5 per session in derivatives/pose/{sub}/{ses}/",
        "Pluggable: SLEAP and LightningPose also supported via movement",
        "Status: 26/26 sessions complete",
    ]
    _add_bullet_box(slide, bullets, font_size=14)
    _add_notes(slide, slide_data["notes"])


def _build_dlc_retrain_slide(prs: Presentation, slide_data: dict):
    """Slide 8: DLC Retraining Workflow -- flow diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "DLC Retraining Workflow")

    inch = 914400
    box_w = int(3.0 * inch)
    box_h = int(0.55 * inch)
    gap = int(0.2 * inch)
    x_center = int(5.15 * inch)

    steps = [
        ("1. Identify Bad Frames", "Tracking QC page -> export frame indices", ACCENT_RED),
        ("2. Extract & Label", "Download video, extract PNGs, label in napari", ACCENT_ORANGE),
        ("3. Upload Labels", "scripts/upload_dlc_labels.py -> S3", MID_BLUE),
        ("4. Train on EC2", "SA fine-tune, up to 50k iters (g4dn/g5)", ACCENT_TEAL),
        ("5. Inference", "Re-run on all 26 sessions", ACCENT_TEAL),
        ("6. Declare Champion", "dlc-champion.json updated on S3", DARK_BLUE),
        ("7. Evaluate", "Compare fine-tuned vs previous in Tracking QC", ACCENT_ORANGE),
        ("8. Downstream", "Re-run Stages 3 -> 3b -> 5 -> 6", DARK_BLUE),
    ]

    for i, (title, desc, colour) in enumerate(steps):
        y = int(1.4 * inch) + i * (box_h + gap)
        _rect(
            slide,
            x_center - box_w // 2,
            y,
            box_w,
            box_h,
            f"{title}\n{desc}",
            fill=colour,
            font_size=11,
        )
        if i < len(steps) - 1:
            _arrow_down(slide, x_center, y + box_h, y + box_h + gap, colour=DARK_GREY)

    # Side labels
    for lbl, y_idx in [
        ("Mac", 1),
        ("Mac", 2),
        ("EC2 GPU", 3),
        ("EC2 GPU", 4),
        ("S3", 5),
        ("Frontend", 6),
        ("Mac/EC2", 7),
    ]:
        y = int(1.4 * inch) + (y_idx - 1) * (box_h + gap)
        lbl_box = slide.shapes.add_textbox(
            Inches(1.0), Emu(int(y + box_h * 0.15)), Inches(1.3), Inches(0.3)
        )
        p = lbl_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        _add_run(p, lbl, size=10, italic=True, colour=MID_GREY)

    _add_notes(slide, slide_data["notes"])


def _build_champion_slide(prs: Presentation, slide_data: dict):
    """Slide 9: DLC Champion Model System."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "DLC Champion Model System")

    bullets = [
        "Single source of truth: s3://hm2p-derivatives/dlc-champion.json",
        "  Contains: champion_id (deterministic hash), model_name, architecture, snapshot",
        "Provenance chain: every DLC-derived file stamps dlc_champion_id as HDF5 attribute",
        "  kinematics.h5 -> sync.h5 -> analysis.h5 (inherited through pipeline)",
        "  Rendered videos get a .provenance.json sidecar",
        "Frontend enforcement:",
        "  get_dlc_champion() -- loads manifest (cached 300s)",
        "  is_session_current() -- compares stored vs current champion_id",
        "  render_champion_staleness_warning() -- shared banner on all DLC-dependent pages",
        "  Stale data is warned, never hidden -- QC must remain possible",
        "History: superseded champions archived to dlc-champion-history/",
    ]
    _add_bullet_box(slide, bullets, font_size=14)
    _add_notes(slide, slide_data["notes"])


def _build_kinematics_slide(prs: Presentation, slide_data: dict):
    """Slide 10: Stage 3 - Kinematics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Stage 3 -- Kinematics")

    # Processing pipeline as numbered list
    steps = [
        "1. Apply per-session orientation rotation (from experiments.csv)",
        "2. Rename SuperAnimal bodypart names to project names",
        "3. Filter low-confidence detections (< 0.9 -> NaN)",
        "4. Interpolate short gaps (up to 5 frames) and smooth",
        "5. Compute HD from ear vector + 3 QC vectors, fuse with confidence weighting",
        "6. Compute position (mm), speed (cm/s), AHV (deg/s), movement state",
        "7. Align light on/off from timestamps.h5; apply bad_behav mask",
        "8. Compute maze coordinates (7x5 grid) via Shapely polygon clipping",
    ]

    bullets = (
        [
            "Input: Pose .h5 from Stage 2 (any tracker format)",
            "Loader: movement.io.load_poses.from_file() -- tracker-agnostic xarray.Dataset",
            "Processing pipeline (kinematics/compute.py):",
        ]
        + ["  " + s for s in steps]
        + [
            "Perspective correction: kinematics/perspective.py -- ground-plane projection",
            "Output: derivatives/movement/{sub}/{ses}/kinematics.h5",
        ]
    )
    _add_bullet_box(slide, bullets, font_size=13)
    _add_notes(slide, slide_data["notes"])


def _build_calcium_slide(prs: Presentation, slide_data: dict):
    """Slide 11: Stage 4 - Calcium Processing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Stage 4 -- Calcium Processing")

    bullets = [
        "Input: Suite2p/CaImAn native files (via roiextractors) + timestamps.h5",
        "Processing pipeline (calcium/run.py):",
        "  1. Neuropil subtraction: F_corr = F - 0.7 * Fneu (fixed) or FISSA (spatial ICA)",
        "  2. Baseline & dF/F0: sliding window minimum of Gaussian-smoothed trace",
        "  3. Event detection: V&H 2020 threshold (primary), SD-threshold (Zong 2022)",
        "  4. CASCADE spike inference: calibrated spikes/s from pre-trained deep-learning models",
        "  5. Per-ROI QC: SNR, decay tau, neuropil correlation, bleach slope, active fraction",
        "Output: derivatives/calcium/{sub}/{ses}/ca.h5",
        "  dff (R, T) float32 -- dF/F0 per ROI",
        "  spikes (R, T) float32 -- CASCADE spike rate (spikes/s)",
        "  event_masks (R, T) float32 -- V&H binary events",
        "  roi_types (R,) uint8, snr (R,) float32",
        "Stage 4b: CASCADE can be re-run independently without repeating neuropil/dF/F",
    ]
    _add_bullet_box(slide, bullets, font_size=13)
    _add_notes(slide, slide_data["notes"])


def _build_sync_slide(prs: Presentation, slide_data: dict):
    """Slide 12: Stage 5 -- Sync -- diagram with shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Stage 5 -- Neural-Behavioural Synchronisation")

    inch = 914400
    input_w = int(2.5 * inch)
    input_h = int(0.7 * inch)

    # Three input boxes at top
    inputs = [
        ("kinematics.h5\n(~100 Hz)", MID_BLUE, 1.0),
        ("ca.h5\n(~9.6 Hz)", ACCENT_TEAL, 5.4),
        ("timestamps.h5\n(frame times)", ACCENT_ORANGE, 9.8),
    ]
    for label, colour, x in inputs:
        _rect(
            slide,
            int(x * inch),
            int(1.5 * inch),
            input_w,
            input_h,
            label,
            fill=colour,
            font_size=13,
        )

    # Central alignment box
    align_w = int(4.0 * inch)
    align_h = int(0.7 * inch)
    align_x = int(4.65 * inch)
    align_y = int(3.0 * inch)
    _rect(
        slide,
        align_x,
        align_y,
        align_w,
        align_h,
        "Align & Resample\n(linear interpolation at imaging frame times)",
        fill=DARK_BLUE,
        font_size=12,
    )

    # Arrows from inputs to alignment
    for x in [1.0, 5.4, 9.8]:
        center_x = int(x * inch) + input_w // 2
        _arrow_down(slide, center_x, int(1.5 * inch) + input_h, align_y, colour=DARK_GREY)

    # Output box
    out_w = int(3.5 * inch)
    out_h = int(0.7 * inch)
    out_x = int(4.9 * inch)
    out_y = int(4.4 * inch)
    _rect(
        slide,
        out_x,
        out_y,
        out_w,
        out_h,
        "sync.h5\n(all signals at ~9.6 Hz imaging rate)",
        fill=ACCENT_GREEN,
        font_size=13,
    )
    _arrow_down(slide, align_x + align_w // 2, align_y + align_h, out_y, colour=DARK_GREY)

    # Details below
    detail_bullets = [
        "Continuous signals (HD, position, speed): linear interpolation",
        "Boolean signals (active, light_on, bad_behav): nearest-neighbour",
        "Calcium data: copied verbatim (already at imaging rate)",
        "DLC provenance: dlc_model_name, dlc_snapshot, dlc_champion_id inherited",
        "Sync diagnostics: per-channel ISI/MAD/CV, cross-channel drift, light protocol validation",
        "Session classification into 7 sync_status tiers (config: sync.yaml)",
    ]
    _add_bullet_box(slide, detail_bullets, top=5.3, height=2.0, font_size=13)
    _add_notes(slide, slide_data["notes"])


def _build_analysis_slide(prs: Presentation, slide_data: dict):
    """Slide 13: Stage 6 - Analysis -- table of modules."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Stage 6 -- Analysis (19 Modules)")

    headers = ["Module", "What it computes"]
    rows = [
        ["tuning.py", "HD tuning curves, preferred direction, mean vector length, Rayleigh test"],
        ["significance.py", "Circular shuffle tests for HD significance"],
        ["classify.py", "Automated HD cell classification"],
        ["comparison.py", "Tuning curve correlation, PD shift, split-half reliability"],
        ["decoder.py", "Bayesian population HD decoder (MAE in degrees)"],
        ["stability.py", "Temporal stability, light/dark drift analysis"],
        ["gain.py", "Light/dark gain modulation index"],
        ["anchoring.py", "Visual vs idiothetic HD anchoring"],
        ["ahv.py", "Angular head velocity tuning"],
        ["speed.py", "Speed modulation analysis"],
        ["information.py", "Spatial/directional information (Skaggs, bits/spike)"],
        ["activity.py", "Active-cell detection, firing rate statistics"],
        ["population.py", "Population-level summary statistics"],
        ["mixed_stats.py", "Cross-module Penk+ vs CamKII+ comparisons"],
        ["celltype_dynamics.py", "Time-resolved population dynamics by cell type"],
        ["rastermap_analysis.py", "Rastermap-based neural population visualisation"],
    ]
    _add_table(
        slide,
        headers,
        rows,
        top=1.3,
        col_widths=[3.0, 9.3],
        row_height=0.32,
        body_font_size=11,
        header_font_size=12,
    )

    # Footer note
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.0), Inches(0.4))
    p = txb.text_frame.paragraphs[0]
    _add_run(
        p,
        "Non-parametric statistics only: Mann-Whitney U, Wilcoxon, Spearman, "
        "Kruskal-Wallis, permutation/bootstrap",
        size=12,
        italic=True,
        colour=MID_GREY,
    )

    _add_notes(slide, slide_data["notes"])


def _build_maze_slide(prs: Presentation, slide_data: dict):
    """Slide 14: Maze Analysis & NaviGraph."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Maze Analysis & NaviGraph")

    bullets = [
        "Maze topology (maze/topology.py): Rosenberg maze as graph "
        "-- 7x5 grid, adjacency, dead-ends",
        "Discretization (maze/discretize.py): continuous (x, y) -> cell assignment via Shapely",
        "Behavioural analysis (maze/analysis.py): occupancy, exploration, "
        "turn bias, dead-end visits",
        "NaviGraph-inspired neural analyses (maze/neural.py):",
        "  1. Light/dark graph annotation: occupancy-normalised activity per ROI per cell",
        "  2. Decision-point HD tuning: split by location type (junction/corridor/dead-end)",
        "  3. Path familiarity: activity change with repeated corridor traversals",
        "  4. Junction choice prediction: cross-validated logistic decoding of turn choice",
        "Citation: Koren Iton A et al. 2025. NaviGraph. bioRxiv. doi:10.1101/2025.05.18.654725",
        "Frontend: maze_page.py (topology + occupancy), "
        "maze_animation_page.py (trajectory playback)",
    ]
    _add_bullet_box(slide, bullets, font_size=14)
    _add_notes(slide, slide_data["notes"])


def _build_canvas_slide(prs: Presentation, slide_data: dict):
    """Slide 15: Canvas Maze Animation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Canvas Maze Animation")

    bullets = [
        "Component: frontend/components/maze_canvas.py -- HTML5 Canvas + JavaScript",
        "Renders at 60 fps without Streamlit reruns (all playback logic in JS)",
        "Visual elements:",
        "  Mouse trajectory trail",
        "  Bodypart skeleton overlay",
        "  Head direction arrow",
        "  Light/dark state indicator",
        "  Maze grid with dead-end shading",
        "Controls: play/pause, speed slider, frame scrubber, zoom",
        "Embedded via st.html(unsafe_allow_javascript=True)",
        "Data serialised as JSON and injected into JS at render time",
    ]
    _add_bullet_box(slide, bullets, font_size=15)
    _add_notes(slide, slide_data["notes"])


def _build_frontend_slide(prs: Presentation, slide_data: dict):
    """Slide 16: Frontend Dashboard -- table of sections."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Frontend Dashboard -- 67 Pages")
    _add_subtitle(slide, "Streamlit with st.navigation() for multi-page layout")

    headers = ["Section", "#", "Key Pages", "Purpose"]
    rows = [
        [
            "Overview",
            "7",
            "Home, Sessions, Animals, Pipeline, Cell Summary",
            "Project status, metadata",
        ],
        [
            "Pipeline",
            "16",
            "Suite2p, Calcium, DLC Training/Inference, Tracking QC, Sync",
            "Per-stage QC diagnostics",
        ],
        [
            "Explore",
            "16",
            "Explorer, Timeline, ROI Gallery, Events, Rastermap, Behaviour",
            "Interactive data exploration",
        ],
        [
            "Analysis",
            "23",
            "HD Tuning, Decoder, Stability, Gain, Anchoring, Maze, Pub Stats",
            "Scientific analysis",
        ],
        ["System", "3", "AWS, Costs, Changelog", "Infrastructure monitoring"],
    ]
    _add_table(
        slide,
        headers,
        rows,
        top=1.6,
        col_widths=[1.8, 0.5, 5.5, 4.5],
        row_height=0.55,
        body_font_size=12,
    )

    bottom_bullets = [
        "Data loading: frontend/data.py -- S3 caching with @st.cache_data",
        "No sidebar filters -- all controls in page body via st.columns()",
        "No synthetic data -- pages load real data from S3, show messages when unavailable",
    ]
    _add_bullet_box(slide, bottom_bullets, top=5.2, height=1.5, font_size=13)
    _add_notes(slide, slide_data["notes"])


def _build_pynapple_slide(prs: Presentation, slide_data: dict):
    """Slide 17: pynapple Interface."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Analysis Loading -- pynapple Interface")

    # Code block as a grey box
    code_lines = [
        "import pynapple as nap, h5py",
        "",
        'with h5py.File("sync.h5") as f:',
        '    t      = f["frame_times"][:]',
        '    spikes = nap.TsdFrame(t=t, d=f["spikes"][:].T)',
        '    dff    = nap.TsdFrame(t=t, d=f["dff"][:].T)',
        '    hd     = nap.Tsd(t=t, d=f["hd_deg"][:])',
        '    speed  = nap.Tsd(t=t, d=f["speed_cm_s"][:])',
    ]
    code_box = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0.5),
        Inches(1.4),
        Inches(8.0),
        Inches(3.0),
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    code_box.line.color.rgb = MID_GREY
    code_box.line.width = Pt(1)

    tf = code_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _add_run(p, line, size=13, colour=DARK_GREY, name=FONT_MONO)

    # Right side: design notes
    design_bullets = [
        "Arrays: time-first, C-contiguous",
        "Timestamps: float64, seconds since session start",
        "Units in names: hd_deg, speed_cm_s, ahv_deg_s",
        "Zero-friction pynapple loading",
    ]
    txb = slide.shapes.add_textbox(Inches(9.0), Inches(1.5), Inches(4.0), Inches(2.5))
    tf2 = txb.text_frame
    tf2.word_wrap = True
    for i, b in enumerate(design_bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_before = Pt(6)
        _add_run(p, b, size=14, colour=DARK_BLUE, bold=True)

    # Planned tools below
    tool_bullets = [
        "Planned tools:",
        "  NEMOS (Flatiron) -- GLM encoding models, pynapple-native, JAX backend",
        "  CEBRA (Schneider et al. 2023) -- contrastive population "
        "embeddings; ring manifold for HD",
        "  neuroconv -- HDF5 -> NWB export for DANDI archiving",
    ]
    _add_bullet_box(slide, tool_bullets, top=4.8, font_size=14)
    _add_notes(slide, slide_data["notes"])


def _build_aws_slide(prs: Presentation, slide_data: dict):
    """Slide 18: AWS Infrastructure -- boxes for services."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Infrastructure -- AWS")

    inch = 914400

    # EC2 instances section
    _rect(
        slide,
        int(0.5 * inch),
        int(1.4 * inch),
        int(5.5 * inch),
        int(2.8 * inch),
        "",
        fill=LIGHT_BLUE,
    )
    txb = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(5.0), Inches(0.35))
    p = txb.text_frame.paragraphs[0]
    _add_run(p, "EC2 Instances (Spot)", size=16, bold=True, colour=DARK_BLUE)

    ec2_types = [
        ("g4dn.xlarge\n~$0.16/hr", ACCENT_TEAL, 0.7, 2.0, 2.3, 0.6),
        ("g4dn.2xl / g5.xl\n~$0.30/hr", ACCENT_TEAL, 3.3, 2.0, 2.5, 0.6),
        ("c5.4xlarge\n~$0.27/hr", MID_BLUE, 0.7, 2.8, 2.3, 0.6),
    ]
    for label, colour, x, y, w, h in ec2_types:
        _rect(
            slide,
            int(x * inch),
            int(y * inch),
            int(w * inch),
            int(h * inch),
            label,
            fill=colour,
            font_size=10,
        )

    # Label the EC2 types
    for label, x, y in [
        ("DLC inference\nSuite2p", 0.7, 2.62),
        ("DLC training", 3.3, 2.62),
        ("CPU stages", 0.7, 3.42),
    ]:
        txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2.3), Inches(0.3))
        p = txb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_run(p, label, size=9, italic=True, colour=MID_GREY)

    # S3 section
    _rect(
        slide,
        int(6.5 * inch),
        int(1.4 * inch),
        int(6.3 * inch),
        int(2.8 * inch),
        "",
        fill=LIGHT_BLUE,
    )
    txb2 = slide.shapes.add_textbox(Inches(6.7), Inches(1.45), Inches(5.8), Inches(0.35))
    p = txb2.text_frame.paragraphs[0]
    _add_run(p, "S3 Buckets", size=16, bold=True, colour=DARK_BLUE)

    _rect(
        slide,
        int(6.8 * inch),
        int(2.0 * inch),
        int(2.7 * inch),
        int(0.6 * inch),
        "hm2p-rawdata\nInfrequent Access",
        fill=ACCENT_ORANGE,
        font_size=11,
    )
    _rect(
        slide,
        int(9.8 * inch),
        int(2.0 * inch),
        int(2.7 * inch),
        int(0.6 * inch),
        "hm2p-derivatives\nStandard",
        fill=ACCENT_GREEN,
        font_size=11,
    )

    # Arrow from EC2 to S3
    _arrow_right(slide, int(6.0 * inch), int(2.8 * inch), int(6.5 * inch))

    # Safety mechanisms below
    safety_bullets = [
        "Safety mechanisms:",
        "  GPU watchdog: terminate if 0% utilisation for 5+ minutes",
        "  24-hour hard timeout on all instances",
        "  Self-termination on completion (InstanceInitiatedShutdownBehavior=terminate)",
        "  Security group restricted to known IP",
        "  SSM Session Manager for keyless SSH",
        "Cost: ~$180-380 one-time processing (all 26 sessions). Storage: ~$10/month",
    ]
    _add_bullet_box(slide, safety_bullets, top=4.5, font_size=13)
    _add_notes(slide, slide_data["notes"])


def _build_orchestration_slide(prs: Presentation, slide_data: dict):
    """Slide 19: Pipeline Orchestration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Pipeline Orchestration")

    bullets = [
        "Snakemake 8.x+ -- DAG-based workflow engine",
        "Snakefile: workflow/Snakefile (main DAG)",
        "Stage rules: workflow/rules/ -- ingest, extraction, pose, kinematics, calcium, sync",
        "Three compute profiles (workflow/profiles/):",
        "  local/ -- CPU-only on laptop",
        "  local-gpu/ -- all stages on local GPU",
        "  aws-batch/ -- full cloud pipeline (managed job queue)",
        "Docker images (docker/):",
        "  gpu.Dockerfile -- CUDA 12.1 + Suite2p + DLC",
        "  cpu.Dockerfile -- CPU-only stages (movement, calcium, sync)",
        "  kpms.Dockerfile -- keypoint-MoSeq isolated environment",
        "  cascade.Dockerfile -- CASCADE with tensorflow (Python 3.8 / TF 2.3)",
        "Direct runner scripts: scripts/run_stage{0,3,4,5,6}_*.py (bypass Snakemake)",
    ]
    _add_bullet_box(slide, bullets, font_size=14)
    _add_notes(slide, slide_data["notes"])


def _build_hdf5_schema_slide(prs: Presentation, slide_data: dict):
    """Slide 20: HDF5 Schemas -- table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Data Standards -- HDF5 Schemas")

    headers = ["File", "Key Datasets", "Shape", "Rate"]
    rows = [
        [
            "timestamps.h5",
            "frame_times_camera, frame_times_imaging, light_on_times",
            "(N,), (T,), (L,)",
            "--",
        ],
        [
            "kinematics.h5",
            "hd_deg, x_mm, y_mm, speed_cm_s, ahv_deg_s, active, light_on",
            "(N,)",
            "~100 Hz",
        ],
        ["ca.h5", "dff, spikes, event_masks, roi_types, snr", "(R,T), (R,), (R,)", "~9.6 Hz"],
        [
            "sync.h5",
            "All kinematics (resampled) + all calcium (verbatim)",
            "(T,), (R,T)",
            "~9.6 Hz",
        ],
        [
            "analysis.h5",
            "tuning_curves, pd, mvl, rayleigh_p, is_hd, decoder_error",
            "(R,B), (R,)",
            "--",
        ],
    ]
    _add_table(
        slide,
        headers,
        rows,
        top=1.3,
        col_widths=[2.5, 5.5, 2.3, 2.0],
        row_height=0.55,
        body_font_size=12,
    )

    design_bullets = [
        "Consistent indexing: arrays are time-first (C-contiguous for fast row slicing)",
        "Timestamps: float64 seconds since session start",
        "Units in dataset names: hd_deg, speed_cm_s, ahv_deg_s, x_mm, y_mm",
        "Validation: pandera schemas in src/hm2p/io/hdf5.py",
        "Provenance attributes: session_id, fps_imaging, dlc_model_name, dlc_champion_id",
        "Backward compatibility aliases: x_mm = x_body_mm, speed_cm_s = speed_body_cm_s",
    ]
    _add_bullet_box(slide, design_bullets, top=4.5, font_size=13)
    _add_notes(slide, slide_data["notes"])


def _build_testing_slide(prs: Presentation, slide_data: dict):
    """Slide 21: Code Quality & Testing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Code Quality & Testing")

    bullets = [
        "117 test files in tests/ mirroring src/hm2p/ structure",
        "Coverage target: >= 90% (hard requirement)",
        "Testing frameworks:",
        "  pytest + pytest-cov -- standard unit tests + coverage",
        "  hypothesis -- property-based testing for numerical functions",
        "  pandera -- runtime schema validation in tests",
        "Pre-commit hooks: ruff (format + lint), mypy (strict), nbstripout, detect-secrets",
        "CI: GitHub Actions -- ci.yml (pytest on Python 3.11 + 3.12), lint.yml (ruff + mypy)",
        "Security: bandit, checkov, detect-secrets, pip-audit, vulture",
        "Test rules:",
        "  Synthetic data only -- never read real data files",
        "  Every function (public and private) must have at least one test",
        "  Non-parametric statistics tests verify parametric alternatives never used",
    ]
    _add_bullet_box(slide, bullets, font_size=14)
    _add_notes(slide, slide_data["notes"])


def _build_source_layout_slide(prs: Presentation, slide_data: dict):
    """Slide 22: Source Code Layout -- nested directory structure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Source Code Layout")

    # Directory structure as monospace in a grey box
    dirs = [
        ("src/hm2p/", ""),
        ("  ingest/", "Stage 0: validate.py, daq.py"),
        ("  extraction/", "Stage 1: suite2p.py, caiman.py, soma_features.py, curation.py"),
        ("  pose/", "Stage 2: run.py, preprocess.py, quality.py, retrain.py, select.py"),
        ("  kinematics/", "Stage 3: compute.py, perspective.py, syllables.py"),
        ("  calcium/", "Stage 4: neuropil.py, dff.py, spikes.py, events.py, qc.py"),
        ("  sync/", "Stage 5: align.py, validate.py, diagnostics.py, report.py"),
        ("  analysis/", "Stage 6: 19 modules (tuning, decoder, stability, ...)"),
        ("  maze/", "Maze topology + discretization + behavioural + neural"),
        ("  anatomy/", "brainreg registration + injection site rendering"),
        ("  patching/", "Patch-clamp electrophysiology pipeline (11 modules)"),
        ("  io/", "HDF5 I/O, S3 paths, NWB export, AWS costs"),
    ]

    code_box = slide.shapes.add_shape(
        1,
        Inches(0.5),
        Inches(1.3),
        Inches(12.3),
        Inches(4.6),
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    code_box.line.color.rgb = MID_GREY
    code_box.line.width = Pt(1)

    tf = code_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)

    for i, (path, desc) in enumerate(dirs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        _add_run(p, path, size=14, bold=True, colour=DARK_BLUE, name=FONT_MONO)
        if desc:
            _add_run(p, f"  -- {desc}", size=13, colour=MID_GREY, name=FONT_MONO)

    footer_bullets = [
        "Total: ~75 Python modules (excluding __init__.py)",
        "Separation of concerns: processing vs analysis vs frontend vs I/O",
        "No circular imports: strict layering (io -> processing -> analysis)",
    ]
    _add_bullet_box(slide, footer_bullets, top=6.1, height=1.2, font_size=13)
    _add_notes(slide, slide_data["notes"])


def _build_patching_slide(prs: Presentation, slide_data: dict):
    """Slide 23: Patching Pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Patching Pipeline")

    bullets = [
        "Purpose: Patch-clamp electrophysiology + morphology for same cell populations",
        "11 modules in src/hm2p/patching/:",
        "  config.py -- pipeline configuration (config/patching.yaml)",
        "  io.py -- WaveSurfer H5 + SWC file I/O",
        "  ephys.py -- electrophysiology signal processing",
        "  protocols.py -- stimulus protocol parsing and response extraction",
        "  spike_features.py -- AP waveform feature extraction",
        "  morphology.py -- SWC morphology loading and analysis",
        "  metrics.py -- intrinsic excitability and passive membrane properties",
        "  statistics.py -- statistical comparisons (non-parametric)",
        "  pca.py -- PCA on electrophysiological features",
        "Frontend: patching_page.py, patching_traces_page.py, patching_morph_page.py",
        "Data: read-only bind mount at /data/patching/",
        "Status: processing modules complete; frontend pages functional",
    ]
    _add_bullet_box(slide, bullets, font_size=14)
    _add_notes(slide, slide_data["notes"])


def _build_status_slide(prs: Presentation, slide_data: dict):
    """Slide 24: Current Status -- two tables."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Current Status")

    # Completed table
    txb1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(3.0), Inches(0.3))
    p = txb1.text_frame.paragraphs[0]
    _add_run(p, "Completed", size=16, bold=True, colour=ACCENT_GREEN)

    completed_headers = ["Component", "Status"]
    completed_rows = [
        ["Stage 0 -- DAQ parsing", "26/26"],
        ["Stage 1 -- Suite2p extraction", "26/26"],
        ["Stage 2 -- DLC pose estimation", "26/26"],
        ["Stage 3 -- Kinematics", "Code complete"],
        ["Stage 4 -- Calcium processing", "26/26"],
        ["Stage 5 -- Sync", "21/21"],
        ["Stage 6 -- Analysis (19 modules)", "Code complete"],
        ["Frontend (67 pages)", "Operational"],
        ["Patching pipeline (11 modules)", "Complete"],
        ["DLC champion model system", "Phases 1-2"],
    ]
    _add_table(
        slide,
        completed_headers,
        completed_rows,
        top=1.45,
        width=6.0,
        col_widths=[4.0, 2.0],
        row_height=0.3,
        body_font_size=11,
        header_font_size=12,
    )

    # Pending table
    txb2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.15), Inches(3.0), Inches(0.3))
    p = txb2.text_frame.paragraphs[0]
    _add_run(p, "Pending", size=16, bold=True, colour=ACCENT_ORANGE)

    pending_headers = ["Component", "Blocker"]
    pending_rows = [
        ["CASCADE spike inference", "TF 2.3 / Python 3.8"],
        ["FISSA neuropil subtraction", "scikit-learn < 1.2"],
        ["NWB export (neuroconv)", "Stub only"],
        ["DLC champion Phase 3", "Partial"],
        ["MoSeq syllables", "Awaiting DLC re-run"],
        ["Snakemake integration", "Scripts used instead"],
    ]
    _add_table(
        slide,
        pending_headers,
        pending_rows,
        left=6.8,
        top=1.45,
        width=6.0,
        col_widths=[3.5, 2.5],
        row_height=0.3,
        body_font_size=11,
        header_font_size=12,
    )

    _add_notes(slide, slide_data["notes"])


def _build_issues_slide(prs: Presentation, slide_data: dict):
    """Slide 25: Known Issues & Documentation Gaps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Known Issues & Documentation Gaps")

    # Left column: doc inconsistencies
    txb1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(5.5), Inches(0.3))
    p = txb1.text_frame.paragraphs[0]
    _add_run(p, "Documentation Inconsistencies", size=14, bold=True, colour=ACCENT_ORANGE)

    doc_bullets = [
        "Count mismatches: README (53 pages), ARCHITECTURE (60), actual (67)",
        "Body part list: PLAN.md lists 5, correct count is 8",
        "8 source modules missing from ARCHITECTURE.md",
        "37 scripts missing from ARCHITECTURE.md",
        "Undocumented features: W&B, NaviGraph, canvas animation, sync diagnostics",
        "README pynapple example uses wrong dataset names",
    ]
    txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(3.5))
    tf = txb2.text_frame
    tf.word_wrap = True
    for i, b in enumerate(doc_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        _add_run(p, b, size=13, colour=DARK_GREY)

    # Right column: technical debt
    txb3 = slide.shapes.add_textbox(Inches(6.8), Inches(1.15), Inches(5.5), Inches(0.3))
    p = txb3.text_frame.paragraphs[0]
    _add_run(p, "Technical Debt", size=14, bold=True, colour=ACCENT_RED)

    tech_bullets = [
        "CASCADE not running: primary spike inference blocked by env constraints",
        "FISSA not running: neuropil subtraction blocked by dependency conflicts",
        "NWB export: stub only -- no data on DANDI",
        "5 sessions missing sync.h5 (likely timing edge cases)",
        "Stale champion model checklist in PLAN.md",
    ]
    txb4 = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.5))
    tf2 = txb4.text_frame
    tf2.word_wrap = True
    for i, b in enumerate(tech_bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_before = Pt(4)
        _add_run(p, b, size=13, colour=DARK_GREY)

    _add_notes(slide, slide_data["notes"])


def _build_reference_slide(prs: Presentation, slide_data: dict):
    """Slide 26: Key File Paths Reference -- table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Key File Paths Reference")

    headers = ["Purpose", "Path"]
    rows = [
        ["Pipeline source code", "src/hm2p/"],
        ["Tests", "tests/"],
        ["Frontend", "frontend/app.py, frontend/data.py, frontend/pages/"],
        ["Pipeline runners", "scripts/run_stage{0,3,4,5,6}_*.py"],
        ["DLC workflow scripts", "scripts/*dlc*.py, scripts/launch_dlc_*.py"],
        ["Config files", "config/pipeline.yaml, compute.yaml, sync.yaml"],
        ["Metadata", "metadata/animals.csv, experiments.csv"],
        ["Snakemake", "workflow/Snakefile, workflow/rules/, workflow/profiles/"],
        ["Docker", "docker/{gpu,cpu,kpms,cascade}.Dockerfile"],
        ["Architecture docs", "ARCHITECTURE.md, PLAN.md, CLAUDE.md"],
        ["Topic docs", "docs/ (41 markdown files)"],
        ["S3 rawdata", "s3://hm2p-rawdata/"],
        ["S3 derivatives", "s3://hm2p-derivatives/"],
        ["DLC champion manifest", "s3://hm2p-derivatives/dlc-champion.json"],
        ["Legacy code (read-only)", "old-pipeline/"],
    ]
    _add_table(
        slide, headers, rows, top=1.3, col_widths=[3.5, 8.8], row_height=0.35, body_font_size=12
    )

    _add_notes(slide, slide_data["notes"])


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

SLIDE_BUILDERS = [
    _build_title_slide,  # 1
    _build_raw_data_slide,  # 2
    _build_data_flow_slide,  # 3
    _build_s3_layout_slide,  # 4
    _build_daq_slide,  # 5
    _build_extraction_slide,  # 6
    _build_dlc_slide,  # 7
    _build_dlc_retrain_slide,  # 8
    _build_champion_slide,  # 9
    _build_kinematics_slide,  # 10
    _build_calcium_slide,  # 11
    _build_sync_slide,  # 12
    _build_analysis_slide,  # 13
    _build_maze_slide,  # 14
    _build_canvas_slide,  # 15
    _build_frontend_slide,  # 16
    _build_pynapple_slide,  # 17
    _build_aws_slide,  # 18
    _build_orchestration_slide,  # 19
    _build_hdf5_schema_slide,  # 20
    _build_testing_slide,  # 21
    _build_source_layout_slide,  # 22
    _build_patching_slide,  # 23
    _build_status_slide,  # 24
    _build_issues_slide,  # 25
    _build_reference_slide,  # 26
]


def main():
    """Build the architecture presentation."""
    repo_root = Path(__file__).resolve().parent.parent
    md_path = repo_root / "docs" / "architecture-slides.md"
    out_path = repo_root / "docs" / "architecture-presentation.pptx"

    if not md_path.exists():
        raise FileNotFoundError(f"Slide deck markdown not found: {md_path}")

    slides_data = _parse_slides(md_path)
    if len(slides_data) != 26:
        raise ValueError(f"Expected 26 slides in markdown, found {len(slides_data)}")

    # Create presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for builder, data in zip(SLIDE_BUILDERS, slides_data, strict=True):
        builder(prs, data)

    prs.save(str(out_path))
    print(f"Saved presentation to {out_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
